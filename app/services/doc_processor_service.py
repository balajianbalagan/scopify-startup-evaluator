
## `app/services/doc_processor_service.py`

import os
import mimetypes
import logging
from typing import Dict, Any
from pptx import Presentation
from google.cloud import documentai_v1 as documentai

from app.core.config import settings

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Config pulled from settings (set these in your env / config)
PROJECT_ID = getattr(settings, "DOCAI_PROJECT_ID", None)
LOCATION = getattr(settings, "DOCAI_LOCATION", "us")
PROCESSOR_ID = getattr(settings, "DOCAI_PROCESSOR_ID", None)


def _docai_client() -> documentai.DocumentProcessorServiceClient:
    return documentai.DocumentProcessorServiceClient()


async def process_file_to_structured(file_path: str) -> Dict[str, Any]:
    """
    Detect file type and return normalized structured output:
    {
      "text": str,
      "tables": list,
      "slides": list,
      "entities": list,
      "metadata": dict
    }
    """
    logger.info(f"Processing file: {file_path}")
    
    if not os.path.exists(file_path):
        logger.error(f"File not found: {file_path}")
        raise FileNotFoundError(f"File not found: {file_path}")

    mime_type, _ = mimetypes.guess_type(file_path)
    logger.info(f"Detected MIME type: {mime_type}")

    if mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        logger.info("Processing as PPTX file")
        return _extract_pptx(file_path)

    # Fallback to Document AI for PDFs and other docs
    logger.info("Processing with Document AI")
    return _extract_docai(file_path)


def _extract_pptx(file_path: str) -> Dict[str, Any]:
    logger.info(f"Extracting content from PPTX: {file_path}")
    try:
        prs = Presentation(file_path)
        logger.info(f"Successfully opened presentation with {len(prs.slides)} slides")
        
        slides_data = []
        all_texts = []

        for i, slide in enumerate(prs.slides, start=1):
            logger.info(f"Processing slide {i}")
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    text = shape.text.strip()
                    slide_texts.append(text)
                    all_texts.append(text)
            slides_data.append({"slide": i, "texts": slide_texts})
            logger.info(f"Slide {i}: Found {len(slide_texts)} text elements")

        logger.info("PPTX extraction completed successfully")
        return {
            "text": "\n".join(all_texts),
            "tables": [],
            "slides": slides_data,
            "entities": [],
            "metadata": {"slides": len(slides_data)}
        }
    except Exception as e:
        logger.error(f"Error extracting PPTX content: {str(e)}", exc_info=True)
        raise


def _extract_docai(file_path: str) -> Dict[str, Any]:
    logger.info(f"Extracting content using Document AI: {file_path}")
    
    if not PROJECT_ID or not PROCESSOR_ID:
        logger.error("DocAI configuration missing")
        raise RuntimeError("DocAI configuration (PROJECT_ID / PROCESSOR_ID) is missing")

    try:
        logger.info("Initializing Document AI client")
        client = _docai_client()
        name = client.processor_path(PROJECT_ID, LOCATION, PROCESSOR_ID)
        logger.info(f"Using processor path: {name}")

        # Guess mime-type; default to application/pdf
        mime_type, _ = mimetypes.guess_type(file_path)
        mime_type = mime_type or "application/pdf"
        logger.info(f"File MIME type: {mime_type}")

        logger.info("Reading file content")
        with open(file_path, "rb") as f:
            raw = f.read()
        logger.info(f"Read {len(raw)} bytes")

        logger.info("Creating Document AI request")
        raw_document = documentai.RawDocument(content=raw, mime_type=mime_type)
        request = documentai.ProcessRequest(name=name, raw_document=raw_document)
        
        logger.info("Processing document with Document AI")
        result = client.process_document(request=request)
        document = result.document
        logger.info("Document processing completed")
        return document
    except Exception as e:
        logger.error(f"Error in Document AI processing: {str(e)}", exc_info=True)
        raise

    text = document.text or ""

    tables = []
    for page in document.pages:
        for table in page.tables:
            rows = []
            for row in table.body_rows:
                cells = []
                for cell in row.cells:
                    # Extract plain text for the cell
                    cell_text = ""
                    if cell.layout and cell.layout.text_anchor:
                        try:
                            for segment in cell.layout.text_anchor.text_segments:
                                start = int(segment.start_index) if segment.start_index else 0
                                end = int(segment.end_index) if segment.end_index else None
                                cell_text += document.text[start:end]
                        except Exception:
                            pass
                    cells.append(cell_text)
                rows.append(cells)
            tables.append(rows)

    entities = []
    try:
        for e in document.entities:
            entities.append({"type": e.type_, "mention_text": e.mention_text})
    except Exception:
        pass

    return {
        "text": text,
        "tables": tables,
        "slides": [],
        "entities": entities,
        "metadata": {"pages": len(document.pages) if document.pages else 0}
    }
